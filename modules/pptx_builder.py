"""
PPTX 빌더 모듈
- 분석된 슬라이드 데이터를 편집 가능한 PowerPoint(.pptx)로 변환
- 텍스트 박스, 이미지, 도형을 정확한 위치에 배치
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import io
import os


# pt를 EMU로 변환 (1pt = 12700 EMU)
PT_TO_EMU = 12700


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """hex 색상 코드를 RGBColor로 변환"""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return RGBColor(0, 0, 0)
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def _alignment(align_str: str):
    """문자열 정렬을 python-pptx 정렬 상수로 변환"""
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }
    return mapping.get(align_str, PP_ALIGN.LEFT)


def build_pptx_from_pdf_data(slides_data: list, output_path: str) -> str:
    """
    PDF에서 직접 추출한 데이터로 PPTX를 생성합니다.

    Args:
        slides_data: pdf_processor.extract_from_pdf()의 결과
        output_path: 출력 PPTX 파일 경로

    Returns:
        str: 생성된 PPTX 파일 경로
    """
    if not slides_data:
        raise ValueError("슬라이드 데이터가 비어있습니다.")

    prs = Presentation()

    # 첫 슬라이드 크기 기준으로 프레젠테이션 크기 설정
    first = slides_data[0]
    prs.slide_width = int(first.width * PT_TO_EMU)
    prs.slide_height = int(first.height * PT_TO_EMU)

    blank_layout = prs.slide_layouts[6]  # 빈 레이아웃

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # 배경색 설정
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(slide_data.background_color)

        # 텍스트 블록 배치
        for tb in slide_data.text_blocks:
            left = int(tb.x * PT_TO_EMU)
            top = int(tb.y * PT_TO_EMU)
            width = int(max(tb.width, 10) * PT_TO_EMU)
            height = int(max(tb.height, 10) * PT_TO_EMU)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = tb.text
            p.font.size = Pt(tb.font_size)
            p.font.color.rgb = _hex_to_rgb(tb.color)
            p.font.bold = tb.bold
            p.font.italic = tb.italic

            # 폰트 이름 설정 (가능한 경우)
            if tb.font_name:
                # 일반적인 PDF 폰트명을 PPTX 호환 이름으로 매핑
                clean_name = tb.font_name.split("+")[-1]  # subset prefix 제거
                clean_name = clean_name.split("-")[0]  # Bold/Italic suffix 제거
                p.font.name = clean_name

        # 이미지 블록 배치
        for ib in slide_data.image_blocks:
            left = int(ib.x * PT_TO_EMU)
            top = int(ib.y * PT_TO_EMU)
            width = int(ib.width * PT_TO_EMU)
            height = int(ib.height * PT_TO_EMU)

            # PIL Image → bytes
            img_stream = io.BytesIO()
            ib.image.save(img_stream, format="PNG")
            img_stream.seek(0)

            slide.shapes.add_picture(img_stream, left, top, width, height)

    prs.save(output_path)
    return output_path


def build_pptx_from_ai_data(layouts: list, page_images: list, output_path: str) -> str:
    """
    AI Vision 분석 결과로 PPTX를 생성합니다.

    Args:
        layouts: ai_analyzer.analyze_slides_batch()의 결과
        page_images: 원본 페이지 이미지 (이미지 요소 및 폴백용)
        output_path: 출력 PPTX 파일 경로

    Returns:
        str: 생성된 PPTX 파일 경로
    """
    if not layouts:
        raise ValueError("레이아웃 데이터가 비어있습니다.")

    prs = Presentation()

    # 16:9 표준 크기 (10" x 5.625")
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    blank_layout = prs.slide_layouts[6]

    for i, layout in enumerate(layouts):
        slide = prs.slides.add_slide(blank_layout)

        # 배경색 설정
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(layout.background_color)

        for elem in layout.elements:
            # 백분율 → EMU 변환
            left = int(slide_w * elem.x / 100)
            top = int(slide_h * elem.y / 100)
            width = int(slide_w * elem.width / 100)
            height = int(slide_h * elem.height / 100)

            # 최소 크기 보장
            width = max(width, Emu(100000))
            height = max(height, Emu(100000))

            if elem.type == "text":
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True

                # 여러 줄 텍스트 처리
                lines = elem.content.split("\n")
                for line_idx, line_text in enumerate(lines):
                    if line_idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = line_text
                    p.font.size = Pt(elem.font_size)
                    p.font.color.rgb = _hex_to_rgb(elem.font_color)
                    p.font.bold = elem.bold
                    p.font.italic = elem.italic
                    p.alignment = _alignment(elem.alignment)

            elif elem.type == "shape":
                # 도형 (사각형) 추가
                shape = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.RECTANGLE
                    left, top, width, height
                )
                if elem.background_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = _hex_to_rgb(elem.background_color)
                shape.line.fill.background()  # 테두리 없음

            elif elem.type == "image":
                # 원본 이미지에서 해당 영역 크롭하여 배치
                if i < len(page_images) and page_images[i] is not None:
                    orig = page_images[i]
                    # 백분율 → 픽셀 좌표
                    crop_x = int(orig.width * elem.x / 100)
                    crop_y = int(orig.height * elem.y / 100)
                    crop_w = int(orig.width * elem.width / 100)
                    crop_h = int(orig.height * elem.height / 100)

                    # 범위 클램핑
                    crop_x = max(0, min(crop_x, orig.width - 1))
                    crop_y = max(0, min(crop_y, orig.height - 1))
                    crop_r = min(crop_x + crop_w, orig.width)
                    crop_b = min(crop_y + crop_h, orig.height)

                    if crop_r > crop_x and crop_b > crop_y:
                        cropped = orig.crop((crop_x, crop_y, crop_r, crop_b))
                        img_stream = io.BytesIO()
                        cropped.save(img_stream, format="PNG")
                        img_stream.seek(0)
                        slide.shapes.add_picture(img_stream, left, top, width, height)

    prs.save(output_path)
    return output_path


def build_pptx_with_background_images(page_images: list, output_path: str) -> str:
    """
    최종 폴백: 각 페이지 이미지를 슬라이드 배경으로 삽입합니다.
    (AI 분석도 실패한 경우의 보험)

    Args:
        page_images: 페이지 이미지 목록
        output_path: 출력 PPTX 파일 경로

    Returns:
        str: 생성된 PPTX 파일 경로
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for img in page_images:
        slide = prs.slides.add_slide(blank_layout)
        img_stream = io.BytesIO()
        img.save(img_stream, format="PNG")
        img_stream.seek(0)

        slide.shapes.add_picture(
            img_stream, 0, 0,
            prs.slide_width, prs.slide_height
        )

    prs.save(output_path)
    return output_path
